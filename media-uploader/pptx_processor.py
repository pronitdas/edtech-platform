import logging
import base64
import io
import re
import hashlib
import os
import json
import time
from typing import Dict, Tuple, List, Any, Optional, Set
from dataclasses import dataclass
from collections import defaultdict
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed

import pptx
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from PIL import Image
from openai import OpenAI

# Reuse TextBlock from pdf_processor.py
@dataclass
class TextBlock:
    """Represents a block of text with its properties."""
    text: str
    font_size: float
    font_name: str
    is_bold: bool
    is_italic: bool
    color: int
    bbox: Tuple[float, float, float, float]
    page_num: int
    block_type: str = "normal"  # Can be: title, heading, paragraph, list_item, table, etc.
    level: int = 0

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTXProcessor:
    """Processor for PPTX files with structure detection similar to PDFProcessor."""
    
    # Default OpenAI API Key - should be loaded from environment in production
    DEFAULT_OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")
    
    # Default model name
    DEFAULT_OPENAI_MODEL = "gpt-4o-mini"
    
    # Batch processing settings
    DEFAULT_BATCH_SIZE = 3
    DEFAULT_MAX_WORKERS = 4
    
    @staticmethod
    def extract_text_blocks(pptx_document: Presentation) -> List[TextBlock]:
        """Extract text blocks with formatting from the presentation."""
        all_blocks = []
        
        # Iterate through slides
        for slide_num, slide in enumerate(pptx_document.slides):
            # Process slide title if present
            if slide.shapes.title and slide.shapes.title.text:
                # Title formatting
                text = slide.shapes.title.text.strip()
                
                # For PPTX, we can approximate the formatting info
                # Title formatting is usually bold and larger
                font_size = 24.0  # Approximate default title size
                font_name = "Default"
                is_bold = True
                is_italic = False
                color = 0  # Default black
                
                # Add slide title
                all_blocks.append(TextBlock(
                    text=text,
                    font_size=font_size,
                    font_name=font_name,
                    is_bold=is_bold,
                    is_italic=is_italic,
                    color=color,
                    bbox=(0, 0, 0, 0),  # Default bounding box
                    page_num=slide_num,
                    block_type="title",
                    level=0
                ))
            
            # Process each shape in the slide
            for shape in slide.shapes:
                # Skip if the shape doesn't have a text frame
                if not hasattr(shape, "text_frame") or not shape.text_frame:
                    continue
                
                # Skip empty text frames
                if not shape.text_frame.text.strip():
                    continue
                    
                # Process each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                        
                    # Determine text level based on indentation
                    level = paragraph.level if hasattr(paragraph, "level") else 0
                    
                    # Extract formatting from first run (approximation)
                    font_size = 12.0  # Default body text size
                    font_name = "Default"
                    is_bold = False
                    is_italic = False
                    color = 0  # Default black
                    
                    # Try to get formatting from runs
                    if paragraph.runs:
                        run = paragraph.runs[0]
                        if hasattr(run, "font"):
                            font = run.font
                            # Font size handling
                            if hasattr(font, "size") and font.size:
                                font_size = font.size.pt if hasattr(font.size, "pt") else font_size
                                
                            # Font name
                            if hasattr(font, "name") and font.name:
                                font_name = font.name
                                
                            # Bold and italic
                            is_bold = bool(font.bold) if hasattr(font, "bold") else False
                            is_italic = bool(font.italic) if hasattr(font, "italic") else False
                            
                            # Color handling
                            if hasattr(font, "color") and hasattr(font.color, "rgb"):
                                # Fix the int() conversion error by handling different types
                                if font.color.rgb:
                                    if isinstance(font.color.rgb, str):
                                        # If it's a string (like '112233'), convert with base 16
                                        color = int(font.color.rgb, 16)
                                    elif isinstance(font.color.rgb, int):
                                        # If it's already an integer
                                        color = font.color.rgb
                                    else:
                                        # Default to 0 (black) for unsupported types
                                        color = 0
                                else:
                                    color = 0
                    
                    # Determine block type
                    block_type = "normal"
                    
                    # Use level info and formatting to determine block type
                    if slide.shapes.title and shape == slide.shapes.title:
                        block_type = "title"
                    elif is_bold and font_size > 14:
                        block_type = "heading"
                    elif paragraph.text.startswith(("•", "-", "*")) or re.match(r'^\d+\.', paragraph.text):
                        block_type = "list_item"
                    elif level > 0:
                        block_type = "list_item"
                    else:
                        block_type = "paragraph"
                    
                    # Create TextBlock
                    all_blocks.append(TextBlock(
                        text=text,
                        font_size=font_size,
                        font_name=font_name,
                        is_bold=is_bold,
                        is_italic=is_italic,
                        color=color,
                        bbox=(0, 0, 0, 0),  # Default bounding box
                        page_num=slide_num,
                        block_type=block_type,
                        level=level
                    ))
        
        return all_blocks
    
    @staticmethod
    def analyze_document_structure(blocks: List[TextBlock]) -> List[TextBlock]:
        """
        Analyze presentation structure based on slide organization.
        Simplified for PPTX as slide organization already provides structure.
        """
        if not blocks:
            return blocks
            
        # Group blocks by slide
        slides = defaultdict(list)
        for block in blocks:
            slides[block.page_num].append(block)
        
        # Process each slide
        for slide_num, slide_blocks in slides.items():
            # Identify slide title as heading level 1
            title_blocks = [b for b in slide_blocks if b.block_type == "title"]
            if title_blocks:
                title_blocks[0].block_type = "heading"
                title_blocks[0].level = 1
            
            # Mark other elements with appropriate levels
            for block in slide_blocks:
                if block.block_type == "normal":
                    if block.is_bold and block.font_size > 14:
                        block.block_type = "heading"
                        block.level = 2
                    elif block.text.startswith(("•", "-", "*")) or re.match(r'^\d+\.', block.text):
                        block.block_type = "list_item"
                        block.level = max(1, block.level)
                    else:
                        block.block_type = "paragraph"
        
        return blocks
    
    @staticmethod
    def organize_blocks_into_document(blocks: List[TextBlock]) -> Dict:
        """
        Organize the classified blocks into a hierarchical document structure.
        Adapted for PPTX with slide-based organization.
        """
        # Group blocks by slide
        slides = defaultdict(list)
        for block in blocks:
            slides[block.page_num].append(block)
        
        # Create document structure
        document = {
            "title": "",
            "content": "",
            "children": []
        }
        
        # Find document title (from first slide title)
        first_slide = slides.get(0, [])
        title_block = next((b for b in first_slide if b.block_type == "heading" and b.level == 1), None)
        if title_block:
            document["title"] = title_block.text
        
        # Process each slide as a section
        for slide_num in sorted(slides.keys()):
            slide_blocks = slides[slide_num]
            
            # Find slide title
            slide_title = next((b.text for b in slide_blocks if b.block_type == "heading" and b.level == 1), 
                              f"Slide {slide_num + 1}")
            
            # Create slide section
            slide_section = {
                "title": slide_title,
                "content": "",
                "children": [],
                "level": 1,
                "type": "slide"
            }
            
            # Add content from non-title blocks
            content_blocks = [b for b in slide_blocks if not (b.block_type == "heading" and b.level == 1)]
            if content_blocks:
                slide_section["content"] = "\n".join([b.text for b in content_blocks])
            
            # Add slide to document
            document["children"].append(slide_section)
        
        return document
    
    @staticmethod
    def extract_images(pptx_document: Presentation) -> Dict[str, Any]:
        """Extract images from a PPTX presentation."""
        images = {}
        image_index = 0
        
        try:
            # First try a more direct approach using presentation-level relationships
            for rel in pptx_document.part.rels.values():
                # Check if relation is an image
                if "image" in rel.reltype:
                    try:
                        # Access image data
                        image_part = rel.target_part
                        image_bytes = image_part.blob
                        
                        if not image_bytes:
                            continue
                            
                        # Process image with PIL
                        image_stream = io.BytesIO(image_bytes)
                        pil_image = Image.open(image_stream)
                        
                        # Get image format
                        image_format = pil_image.format.lower() if pil_image.format else "png"
                        
                        # Create unique ID
                        image_hash = hashlib.md5(image_bytes).hexdigest()[:10]
                        image_filename = f"img_pres_{image_index}_{image_hash}.{image_format}"
                        
                        # Skip small images
                        if pil_image.width < 20 or pil_image.height < 20:
                            continue
                            
                        # Convert to base64
                        img_b64 = base64.b64encode(image_bytes).decode()
                        
                        # Store image
                        images[image_filename] = {
                            "data": img_b64,
                            "format": image_format,
                            "width": pil_image.width,
                            "height": pil_image.height,
                            "page": 0,  # Presentation-level image
                            "mode": pil_image.mode,
                            "hash": image_hash
                        }
                        
                        image_index += 1
                    except Exception as e:
                        logger.warning(f"Error extracting presentation-level image: {str(e)}")
        except Exception as e:
            logger.warning(f"Error processing presentation-level images: {str(e)}")
            
        # Process each slide for images
        for slide_num, slide in enumerate(pptx_document.slides):
            try:
                # Try to access slide-level relationships
                for rel in slide.part.rels.values():
                    if "image" in rel.reltype:
                        try:
                            # Access image part data
                            image_part = rel.target_part
                            image_bytes = image_part.blob
                            
                            if not image_bytes:
                                continue
                                
                            # Process image with PIL
                            image_stream = io.BytesIO(image_bytes)
                            pil_image = Image.open(image_stream)
                            
                            # Get image format
                            image_format = pil_image.format.lower() if pil_image.format else "png"
                            
                            # Create unique ID
                            image_hash = hashlib.md5(image_bytes).hexdigest()[:10]
                            image_filename = f"img_slide{slide_num+1}_{image_index}_{image_hash}.{image_format}"
                            
                            # Skip small images
                            if pil_image.width < 20 or pil_image.height < 20:
                                continue
                                
                            # Convert to base64
                            img_b64 = base64.b64encode(image_bytes).decode()
                            
                            # Store image
                            images[image_filename] = {
                                "data": img_b64,
                                "format": image_format,
                                "width": pil_image.width,
                                "height": pil_image.height,
                                "page": slide_num + 1,
                                "mode": pil_image.mode,
                                "hash": image_hash
                            }
                            
                            image_index += 1
                        except Exception as e:
                            logger.warning(f"Error extracting image from slide {slide_num + 1}: {str(e)}")
            except Exception as e:
                logger.warning(f"Error processing slide {slide_num + 1} images: {str(e)}")
            
            # Also try the shape-based approach as fallback
            for shape in slide.shapes:
                # Check if shape is a picture
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        # Try to access the image through shape properties
                        image = shape._pic
                        image_rid = image.xpath('./p:blipFill/a:blip/@r:embed')[0]
                        
                        # Get the relationship
                        if hasattr(slide.part, 'related_parts'):
                            image_part = slide.part.related_parts[image_rid]
                        elif hasattr(slide.part, 'rels'):
                            image_part = slide.part.rels[image_rid].target_part
                        else:
                            continue
                            
                        image_bytes = image_part.blob
                        
                        # Skip if no image data
                        if not image_bytes:
                            continue
                        
                        # Create a unique hash to avoid duplicates with those already extracted
                        image_hash = hashlib.md5(image_bytes).hexdigest()[:10]
                        
                        # Check if we already extracted this image
                        duplicate = False
                        for existing in images.values():
                            if existing.get("hash") == image_hash:
                                duplicate = True
                                break
                                
                        if duplicate:
                            continue
                        
                        # Process the image
                        image_stream = io.BytesIO(image_bytes)
                        pil_image = Image.open(image_stream)
                        
                        # Determine image format
                        image_format = pil_image.format.lower() if pil_image.format else "png"
                        
                        image_filename = f"img_shape{slide_num+1}_{image_index}_{image_hash}.{image_format}"
                        
                        # Skip small images
                        if pil_image.width < 20 or pil_image.height < 20:
                            continue
                            
                        # Convert to base64
                        img_b64 = base64.b64encode(image_bytes).decode()
                        
                        # Store image
                        images[image_filename] = {
                            "data": img_b64,
                            "format": image_format,
                            "width": pil_image.width,
                            "height": pil_image.height,
                            "page": slide_num + 1,
                            "mode": pil_image.mode,
                            "hash": image_hash
                        }
                        
                        image_index += 1
                    except Exception as e:
                        # Just log as debug since we're using multiple fallback approaches
                        logger.debug(f"Error with shape-based image extraction on slide {slide_num + 1}: {str(e)}")
        
        return images
    
    @staticmethod
    def document_to_text(document: Dict) -> str:
        """
        Convert document structure to flat text with markdown formatting.
        Adapted for PPTX to handle slides as sections.
        """
        parts = []
        
        # Add title
        if document.get("title"):
            parts.append(f"# {document['title']}\n")
        
        # Add content
        if document.get("content"):
            parts.append(document["content"])
        
        # Process slides as sections
        for slide in document.get("children", []):
            # Add slide title with heading level 2
            parts.append(f"\n## {slide['title']}\n")
            
            # Add slide content
            if slide.get("content"):
                parts.append(slide["content"])
        
        return "\n".join(parts)
    
    @staticmethod
    def process_pptx(file_data: bytes) -> Tuple[str, Dict[str, Any], Dict[str, Any]]:
        """
        Process a PPTX file with structure detection.
        Returns: (text_content, images_dict, metadata_dict)
        """
        try:
            # Load the presentation from bytes
            pptx_stream = io.BytesIO(file_data)
            pptx_document = Presentation(pptx_stream)
            
            # Extract basic metadata
            slide_count = len(pptx_document.slides)
            metadata = {
                "title": "",  # PPTX doesn't have direct metadata access like DOCX
                "format": "PPTX",
                "pages": slide_count,
                "file_size": len(file_data),
            }
            
            # Extract rich text blocks with formatting
            text_blocks = PPTXProcessor.extract_text_blocks(pptx_document)
            
            # Update metadata with presentation title if found
            title_block = next((b for b in text_blocks 
                              if b.page_num == 0 and b.block_type == "title"), None)
            if title_block:
                metadata["title"] = title_block.text
            
            # Analyze document structure
            classified_blocks = PPTXProcessor.analyze_document_structure(text_blocks)
            
            # Organize into hierarchical document
            document_structure = PPTXProcessor.organize_blocks_into_document(classified_blocks)
            
            # Extract images
            images = PPTXProcessor.extract_images(pptx_document)
            
            # Generate flat text representation
            text_content = PPTXProcessor.document_to_text(document_structure)
            
            return text_content, images, metadata
            
        except Exception as e:
            logger.error(f"Error converting PPTX: {str(e)}")
            # Fallback to basic text extraction
            try:
                pptx_stream = io.BytesIO(file_data)
                pptx_document = Presentation(pptx_stream)
                
                text_content = []
                for slide_idx, slide in enumerate(pptx_document.slides):
                    text_content.append(f"Slide {slide_idx + 1}")
                    
                    # Extract text from all shapes in the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            text = shape.text_frame.text.strip()
                            if text:
                                text_content.append(text)
                    
                    # Add a separator between slides
                    text_content.append("")
                
                return "\n".join(text_content), {}, {"format": "PPTX", "pages": len(pptx_document.slides)}
            except:
                raise  # Re-raise if even the fallback fails
    
    # Reuse methods from PDFProcessor for textbook conversion
    @staticmethod
    def parse_to_textbook(text: str) -> Dict:
        """Reuse the PDF textbook parsing logic for consistent output structure."""
        from pdf_processor import PDFProcessor
        return PDFProcessor.parse_pdf_to_textbook(text)
    
    @staticmethod
    def convert_tree_to_textbook(tree: Dict, knowledge_id: int, knowledge_name: str) -> Dict:
        """Reuse the PDF textbook conversion logic."""
        from pdf_processor import PDFProcessor
        return PDFProcessor.convert_tree_to_textbook(tree, knowledge_id, knowledge_name)
    
    @staticmethod
    def walk_textbook(textbook: Dict, knowledge_id: int) -> List[Dict]:
        """Reuse the PDF textbook walking logic."""
        from pdf_processor import PDFProcessor
        return PDFProcessor.walk_textbook(textbook, knowledge_id)
    
    @staticmethod
    def process_text_to_index(text: str, knowledge_id: int, knowledge_name: str,
                             openai_api_key: str = DEFAULT_OPENAI_API_KEY,
                             openai_model: str = DEFAULT_OPENAI_MODEL,
                             batch_size: int = DEFAULT_BATCH_SIZE,
                             max_workers: int = DEFAULT_MAX_WORKERS) -> Tuple[Dict, List[Dict]]:
        """
        Process presentation text into structured format using OpenAI API.
        
        Directly uses VideoProcessorV2 methods:
        1. Split the text into manageable chunks
        2. Process each chunk with OpenAI to generate structured content
        3. Merge the chunks into a coherent course structure
        4. Convert the structure to chapters for database insertion
        
        Returns (course_structure, chapters)
        """
        logger.info(f"Processing PPTX text with smart chapter generation for knowledge ID {knowledge_id}")
        
        # Import VideoProcessorV2 directly
        from video_processor_v2 import VideoProcessorV2
        
        # Initialize OpenAI client
        client = OpenAI(api_key=openai_api_key)
        
        # Step 1: Split text into chunks using VideoProcessorV2
        chunks = VideoProcessorV2.chunk_text(text)
        logger.info(f"Split PPTX text into {len(chunks)} chunks")
        
        # Step 2: Process chunks in batches
        # Define a function that adapts VideoProcessorV2's process_chunk method for PPTX
        def process_pptx_chunk(chunk, chunk_index, total_chunks, client, model_name):
            # Use the same schema but adapt the prompt for PPTX content
            try:
                prompt = f"""
                Transform this presentation content into a structured course chapter, adding useful insights or examples where appropriate.
                
                PRESENTATION CONTENT:
                ```
                {chunk}
                ```
                
                INSTRUCTIONS:
                1. Create a coherent chapter with a clear title and sections based on the content
                2. Preserve all original content and examples
                3. Add appropriate insights, examples, or clarifications where helpful
                4. Include all numerical examples and data points from the original
                5. This is part {chunk_index+1} of {total_chunks}, ensure your section flows with others
                
                Structure the content with proper headings, key points, and examples. Make it thorough and educational.
                """

                # Use VideoProcessorV2's process method but with our PPTX-specific prompt
                response = client.chat.completions.create(
                    model=model_name,
                    messages=[
                        {
                            "role": "system",
                            "content": "You are a curriculum development expert specializing in creating educational content from presentations.",
                        },
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.5,
                    response_format={"type": "json_schema", "json_schema": VideoProcessorV2.get_chapter_schema()},
                )

                # Extract the structured content
                structured_content = json.loads(response.choices[0].message.content)

                # Add chunk metadata
                structured_content["chunk_index"] = chunk_index
                structured_content["total_chunks"] = total_chunks
                
                # Add PPTX-specific metadata
                structured_content["document_type"] = "pptx"

                return structured_content
                
            except Exception as e:
                logger.error(f"Error processing chunk {chunk_index+1}: {str(e)}")
                # Return an error structure to maintain consistency
                return {
                    "title": f"Error in Chunk {chunk_index+1}",
                    "sections": [
                        {
                            "heading": "Processing Error",
                            "content": f"An error occurred while processing this chunk: {str(e)}",
                            "key_points": [],
                            "examples": [],
                        }
                    ],
                    "chapter_number": chunk_index + 1,
                    "learning_objectives": ["Resolve processing errors"],
                    "chunk_index": chunk_index,
                    "total_chunks": total_chunks,
                    "error": str(e),
                    "document_type": "pptx"
                }
        
        # Process chunks with our custom function
        chunk_results = []
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = [
                executor.submit(
                    process_pptx_chunk,
                    chunk,
                    i,
                    len(chunks),
                    client,
                    openai_model,
                )
                for i, chunk in enumerate(chunks)
            ]
            
            # Process results as they complete
            for i, future in enumerate(as_completed(futures)):
                try:
                    result = future.result()
                    chunk_results.append(result)
                    logger.info(f"Completed chunk {i+1}/{len(chunks)}")
                    
                    # Add a small delay between batches to avoid rate limiting
                    if (i + 1) % batch_size == 0:
                        time.sleep(1)
                except Exception as e:
                    logger.error(f"Error in future for chunk {i}: {str(e)}")
                    chunk_results.append({
                        "title": f"Error in Chunk Processing",
                        "sections": [
                            {
                                "heading": "Processing Error",
                                "content": f"An error occurred: {str(e)}",
                                "key_points": [],
                                "examples": [],
                            }
                        ],
                        "chapter_number": i + 1,
                        "learning_objectives": [],
                        "chunk_index": i,
                        "total_chunks": len(chunks),
                        "error": str(e),
                        "document_type": "pptx"
                    })
        
        # Sort results by chunk index to maintain order
        chunk_results.sort(key=lambda x: x.get("chunk_index", 0))
        
        # Step 3: Merge chunks into course structure using VideoProcessorV2
        course_structure = VideoProcessorV2.merge_structured_chunks(
            chunk_results=chunk_results,
            knowledge_name=knowledge_name,
            client=client,
            model_name=openai_model,
        )
        
        # Add PPTX-specific metadata
        course_structure["metadata"] = {
            "markdown": "",
            "metadata": {
                "processed_at": datetime.utcnow().isoformat(),
                "document_type": "pptx"
            }
        }
        
        # Step 4: Create chapters from structure using VideoProcessorV2
        chapters = VideoProcessorV2.create_chapters_from_structure(
            course_structure=course_structure, 
            knowledge_id=knowledge_id
        )
        
        return course_structure, chapters
    
    @staticmethod
    def prepare_images_for_upload(images: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
        """
        Prepare extracted images for upload - custom implementation for PPTX.
        
        Args:
            images: Dictionary of images extracted from PPTX
            
        Returns:
            Dictionary ready for upload with standardized metadata and buffer
        """
        upload_ready = {}
        
        for filename, img_data in images.items():
            # Extract the base64 encoded data
            img_base64 = img_data.get("data", "")
            
            # Convert base64 to bytes for upload
            img_bytes = base64.b64decode(img_base64) if img_base64 else b""
            
            # Get image format
            img_format = img_data.get("format", "png")
            
            # Create standardized metadata
            upload_ready[filename] = {
                "buffer": img_bytes,  # Add the bytes data as buffer
                "format": img_format,
                "width": img_data.get("width", 0),
                "height": img_data.get("height", 0),
                "page": img_data.get("page", 1),
                "metadata": {
                    "filename": filename,
                    "format": img_format,
                    "width": img_data.get("width", 0),
                    "height": img_data.get("height", 0),
                    "page": img_data.get("page", 1),
                    "mode": img_data.get("mode", ""),
                    "hash": img_data.get("hash", ""),
                    "size_bytes": len(img_bytes) if img_bytes else 0
                }
            }
        
        return upload_ready 